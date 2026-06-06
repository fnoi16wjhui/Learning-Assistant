"""Office document material extraction."""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

from PIL import Image

from src.materials.extractors.base import MaterialExtractor, MaterialParseError
from src.materials.extractors.image import ocr_pil_image
from src.materials.models import MaterialSegment, MaterialType


class DocxExtractor(MaterialExtractor):
    material_type = MaterialType.DOCX
    suffixes = {".docx"}

    def extract(self, path: Path) -> list[MaterialSegment]:
        try:
            from docx import Document
        except ImportError as exc:
            raise MaterialParseError("DOCX parsing requires python-docx. Run: pip install python-docx") from exc

        try:
            document = Document(str(path))
        except Exception as exc:
            raise MaterialParseError(f"DOCX open failed: path={path}") from exc

        parts: list[str] = []
        for paragraph in document.paragraphs:
            if paragraph.text.strip():
                parts.append(paragraph.text)
        for table in document.tables:
            for row in table.rows:
                values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if values:
                    parts.append(" | ".join(values))
        # OCR embedded images inside the DOCX archive
        try:
            with zipfile.ZipFile(path) as z:
                for name in z.namelist():
                    if not name.startswith("word/media/"):
                        continue
                    try:
                        image_bytes = z.read(name)
                        pil_image = Image.open(io.BytesIO(image_bytes))
                        ocr_text = ocr_pil_image(pil_image)
                        if ocr_text.strip():
                            parts.append(ocr_text.strip())
                    except Exception:
                        pass
        except Exception:
            pass
        return [MaterialSegment(text="\n".join(parts))] if parts else []


class PptxExtractor(MaterialExtractor):
    material_type = MaterialType.PPTX
    suffixes = {".pptx"}

    def extract(self, path: Path) -> list[MaterialSegment]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise MaterialParseError("PPTX parsing requires python-pptx. Run: pip install python-pptx") from exc

        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise MaterialParseError(f"PPTX open failed: path={path}") from exc

        segments: list[MaterialSegment] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    texts.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if values:
                            texts.append(" | ".join(values))
            # OCR embedded images in shapes
            for shape in slide.shapes:
                try:
                    image_blob = shape.image.blob
                    pil_image = Image.open(io.BytesIO(image_blob))
                    ocr_text = ocr_pil_image(pil_image)
                    if ocr_text.strip():
                        texts.append(ocr_text.strip())
                except Exception:
                    pass
            if texts:
                segments.append(MaterialSegment(text="\n".join(texts), slide=index))
        return segments
